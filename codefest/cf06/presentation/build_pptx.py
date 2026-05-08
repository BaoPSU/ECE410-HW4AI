from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT  = "/home/bao/ECE410-HW4AI/codefest/cf06/presentation/"
PPTX = OUT + "cf06_cllm_slides.pptx"

BLUE  = RGBColor(0x1f, 0x4e, 0x79)
GREEN = RGBColor(0x1a, 0x7a, 0x4a)
BLACK = RGBColor(0x11, 0x11, 0x11)
WHITE = RGBColor(0xff, 0xff, 0xff)

W, H = Inches(13.33), Inches(7.5)
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(blank)

def section_label(sl, text):
    tb = sl.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(12.0), Inches(0.45))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = GREEN

def script_text(sl, paragraphs, size=26):
    tb = sl.shapes.add_textbox(Inches(0.4), Inches(0.65), Inches(12.5), Inches(6.7))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(10) if para == "" else Pt(0)
        r = p.add_run()
        r.text = para
        r.font.size = Pt(size)
        r.font.color.rgb = BLACK

# ── SLIDE 1: TITLE ────────────────────────────────────────────────────────────
sl = add_slide()
s = sl.shapes.add_shape(1, 0, 0, W, H)
s.line.fill.background(); s.fill.solid(); s.fill.fore_color.rgb = BLUE

tb = sl.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.8))
tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Codefest 6 — CLLM"
r.font.size = Pt(60); r.font.bold = True; r.font.color.rgb = WHITE

tb2 = sl.shapes.add_textbox(Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.8))
tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "4×4 Binary-Weight Crossbar MAC  ·  Claude Sonnet 4.6"
r2.font.size = Pt(24); r2.font.color.rgb = RGBColor(0xbf, 0xd7, 0xed)

# ── SLIDE 2: THE DESIGN — PART 1 ─────────────────────────────────────────────
sl = add_slide()
section_label(sl, "The Design")
script_text(sl, [
    "For CF06 CLLM I had Claude Sonnet 4.6 generate a 4-by-4 binary-weight crossbar MAC unit in SystemVerilog.",
    "",
    "So about the design. The module takes four 8-bit signed inputs, a 4-by-4 weight matrix where each weight is either plus-one or minus-one, and produces four 10-bit signed outputs.",
    "",
    "Each clock cycle, every output computes a dot product, so basically you multiply each input by its weight at that row-column intersection and add them up.",
    "",
    "That's what the crossbar does, a grid of wires where rows carry inputs, columns carry outputs, and the weights sit at the intersections.",
])

# ── SLIDE 3: THE DESIGN — PART 2 (3 STAGES) ──────────────────────────────────
sl = add_slide()
section_label(sl, "The Design — 3 Stages")
script_text(sl, [
    "The module has three stages. First stage is the weight register, it's a flip-flop that latches the full weight matrix when you pulse weight_load high.",
    "",
    "Claude packed the entire 4-by-4 matrix into a single 16-bit flat register, where bit 4i plus j holds weight[i][j]. For an example, bit zero is weight[0][0], bit four would be weight[1][0], and so on.",
    "",
    "I asked Claude why it went with a flat register and it said it actually loads the entire matrix in exactly one clock cycle. Like if you loaded row by row, that would've taken four cycles and you'd need an address counter to track which row you're on, so the flat register just keeps it simple.",
])

# ── SLIDE 4: THE DESIGN — PART 3 (MAC + BIT WIDTH) ───────────────────────────
sl = add_slide()
section_label(sl, "The Design — MAC & Bit Width")
script_text(sl, [
    "The second stage is the combinational MAC, those are just wires, no clock, which sign-extends the inputs to 10 bits and computes all four dot products at the same time.",
    "",
    "Then the third stage is the output register, which latches the results on the next rising edge.",
    "",
    "As for the output bit width, worst case is four inputs at plus or minus 127, so 4 times 127 is 508. And 10-bit signed goes up to 511, which clears the threshold we need.",
])

# ── SLIDE 5: THE TESTBENCH — PART 1 ──────────────────────────────────────────
sl = add_slide()
section_label(sl, "The Testbench")
script_text(sl, [
    "So about the testbench. Claude loaded the weight matrix from the assignment spec, which is plus-one, minus-one, plus-one, minus-one on row zero, plus-one, plus-one, minus-one, minus-one on row one, and so on, and then it applied inputs of 10, 20, 30, and 40.",
    "",
    "Before running the simulation I hand-calculated the expected outputs, so for example, column zero, you get plus-ten from row 0, plus-twenty from row 1, and then minus-thirty and minus-forty from rows 2 and 3, and so that adds up to minus-40.",
    "",
    "I did the same for all four columns and got minus-40, zero, minus-20, and minus-20.",
])

# ── SLIDE 6: THE TESTBENCH — PART 2 (TIMING) ─────────────────────────────────
sl = add_slide()
section_label(sl, "The Testbench — Timing")
script_text(sl, [
    "One thing about the timing Claude had to get right, and the output shows up two cycles after you assert weight_load, not one.",
    "",
    "Basically here's why, so on the first rising edge after weight_load, the weight register latches the new weights, but the output register clocks on that same edge and it's still using the old weights.",
    "",
    "And so the correct output only appears one cycle later. The testbench explicitly waits two cycles before reading the outputs.",
])

# ── SLIDE 7: THE SIMULATION RESULTS ──────────────────────────────────────────
sl = add_slide()
section_label(sl, "The Simulation Results")
script_text(sl, [
    "All four outputs match what I expected.",
    "",
    "Out-zero is minus-40, out-one is zero, out-two and out-three are both minus-20, exactly what my hand calculation said.",
    "",
    "So 4 out of 4, that's pretty good. Thanks.",
], size=32)

prs.save(PPTX)
print(f"Saved — {len(prs.slides)} slides")
